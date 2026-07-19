"""Agrega el link del sitio web del proyecto a la presentación (PPTX).

Inserta un cuadro de texto con el enlace clicable en la diapositiva 12
("Presentación de producto o prototipo"), justo debajo del texto existente.

Uso:
    python3 agregar_link_ppt.py
Modifica: "Presentación Proyecto.pptx" (en el mismo directorio)
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

ARCHIVO = os.path.join(os.path.dirname(__file__), "..", "Presentación Proyecto.pptx")
URL_SITIO = "https://proyecto-calculo-optimizacion-cloud.vercel.app"
SLIDE_PRODUCTO = 11  # índice 0 -> Diapositiva 12: Presentación de producto o prototipo

prs = Presentation(ARCHIVO)
slide = prs.slides[SLIDE_PRODUCTO]

# Ubicar el cuadro de texto existente para posicionar el nuevo justo debajo
ref = slide.shapes[0]
left = ref.left
top = ref.top + ref.height + Emu(100000)
width = ref.width
height = Emu(500000)

caja = slide.shapes.add_textbox(left, top, width, height)
tf = caja.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run_etiqueta = p.add_run()
run_etiqueta.text = "Sitio web del proyecto (incluye simulador interactivo): "
run_etiqueta.font.size = Pt(14)
run_etiqueta.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

run_link = p.add_run()
run_link.text = URL_SITIO
run_link.font.size = Pt(14)
run_link.font.bold = True
run_link.font.color.rgb = RGBColor(0x22, 0x63, 0xEB)
run_link.hyperlink.address = URL_SITIO

prs.save(ARCHIVO)
print(f"Link agregado en la diapositiva {SLIDE_PRODUCTO + 1} de {ARCHIVO}")
