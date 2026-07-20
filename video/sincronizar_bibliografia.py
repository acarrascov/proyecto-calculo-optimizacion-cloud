"""Sincroniza la bibliografía del PPT (diapositiva 14) con la del sitio web
(web/index.html), que quedó con solo 3 fuentes tras quitar la referencia de
Stewart. Reutiliza el formato de viñeta (buChar/buClr) ya aplicado por
estilizar_ppt.py, clonando el primer párrafo de bibliografía como plantilla.

Uso:
    python3 sincronizar_bibliografia.py
"""
import os
import copy
from pptx import Presentation

ARCHIVO = os.path.join(os.path.dirname(__file__), "..", "Presentación Proyecto.pptx")

# Debe coincidir exactamente con la lista de web/index.html (#bibliografia)
REFERENCIAS = [
    "Documentación oficial de precios de AWS EC2: https://aws.amazon.com/ec2/pricing/",
    "Documentación de sympy: https://docs.sympy.org/",
    "Apuntes de clases — Unidad 3: Aplicación de la Derivada.",
]


def main():
    prs = Presentation(ARCHIVO)
    slide = prs.slides[13]  # Diapositiva 14: Bibliografía

    caja = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "CuadroTexto 3":
            caja = shape
            break
    if caja is None:
        raise RuntimeError("No se encontró el cuadro de texto de bibliografía")

    tf = caja.text_frame
    txBody = tf._txBody
    parrafos = tf.paragraphs

    # Párrafo 2 (índice 2) es el primer bullet de referencia; se usa como
    # plantilla de formato (bullet, tamaño, color, fuente).
    plantilla = parrafos[2]._p

    # Quitar todos los párrafos de referencias existentes (índices 2 en adelante)
    for p in parrafos[2:]:
        txBody.remove(p._p)

    for texto in REFERENCIAS:
        nuevo_p = copy.deepcopy(plantilla)
        # Dejar un solo run con el texto nuevo, conservando el rPr del primer run
        runs_r = nuevo_p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
        )
        primer_r = runs_r[0]
        for r in runs_r[1:]:
            nuevo_p.remove(r)
        t_el = primer_r.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
        )
        t_el.text = texto
        txBody.append(nuevo_p)

    prs.save(ARCHIVO)
    print(f"Bibliografía sincronizada ({len(REFERENCIAS)} fuentes) en {ARCHIVO}")


if __name__ == "__main__":
    main()
