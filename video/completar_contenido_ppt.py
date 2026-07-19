"""Completa el contenido real de las 14 diapositivas de "Presentación Proyecto.pptx",
respetando EXACTAMENTE la estructura/orden de temas que define la plantilla original
(cada diapositiva ya indicaba, a modo de instrucción, qué debía contener).

Importante: este script asume que "agregar_link_ppt.py" ya se ejecutó antes (el
link del sitio ya está en la diapositiva 12). No lo vuelve a tocar, solo agrega
contenido nuevo alrededor de él.

Uso:
    python3 completar_contenido_ppt.py
Modifica: "Presentación Proyecto.pptx" (en el directorio padre de video/)
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

ARCHIVO = os.path.join(os.path.dirname(__file__), "..", "Presentación Proyecto.pptx")
URL_SITIO = "https://proyecto-calculo-optimizacion-cloud.vercel.app"

NEGRO = RGBColor(0x22, 0x22, 0x22)
GRIS = RGBColor(0x55, 0x55, 0x55)

# Márgenes consistentes con los títulos originales de la plantilla
IZQ = Emu(838200)
ANCHO = Emu(10515600)
TOPE_INFERIOR = Emu(6558000)  # deja ~300000 EMU de margen inferior


def limpiar_y_titular(text_frame, titulo, tam_titulo=26, negrita=True, alinear=None):
    """Vacía el text frame y escribe el título como primer párrafo."""
    text_frame.clear()
    text_frame.word_wrap = True
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    p = text_frame.paragraphs[0]
    if alinear:
        p.alignment = alinear
    run = p.add_run()
    run.text = titulo
    run.font.size = Pt(tam_titulo)
    run.font.bold = negrita
    run.font.name = "Calibri"
    run.font.color.rgb = NEGRO
    return text_frame


def agregar_parrafo(text_frame, texto, tam=18, negrita=False, color=NEGRO, espacio_antes=None):
    p = text_frame.add_paragraph()
    if espacio_antes is not None:
        p.space_before = Pt(espacio_antes)
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tam)
    run.font.bold = negrita
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return p


def agregar_bullets(text_frame, items, tam=18, color=NEGRO):
    for item in items:
        agregar_parrafo(text_frame, f"•  {item}", tam=tam, color=color)


def reposicionar(shape, left=IZQ, top=Emu(365125), width=ANCHO, height=None):
    shape.left = left
    shape.top = top
    shape.width = width
    if height is not None:
        shape.height = height


def main():
    prs = Presentation(ARCHIVO)
    slides = prs.slides

    # ---------- Diapositiva 1: Portada ----------
    s = slides[0]
    titulo_shape = s.shapes[0]
    tf = limpiar_y_titular(
        titulo_shape.text_frame,
        "Minimización del costo de alojamiento en la nube",
        tam_titulo=34,
        alinear=PP_ALIGN.CENTER,
    )
    # Caja adicional con datos del curso, debajo del título
    top_info = Emu(titulo_shape.top + titulo_shape.height + 100000)
    caja_info = s.shapes.add_textbox(titulo_shape.left, top_info, titulo_shape.width, Emu(1800000))
    tf_info = caja_info.text_frame
    tf_info.word_wrap = True
    p0 = tf_info.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    for i, linea in enumerate([
        "Andrés Carrasco Valdés",
        "Profesor: Carlos Castro Maldonado",
        "Asignatura: Cálculo Diferencial",
        "Carrera: Ingeniería en Informática",
    ]):
        p = p0 if i == 0 else tf_info.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = linea
        run.font.size = Pt(20)
        run.font.name = "Calibri"
        run.font.color.rgb = GRIS

    # ---------- Diapositiva 2: Definición del proyecto ----------
    s = slides[1]
    tf = limpiar_y_titular(s.shapes[0].text_frame, "Definición del proyecto")
    reposicionar(s.shapes[0], height=Emu(1325563))
    agregar_bullets(tf, [
        "Proyecto de Aprendizaje Basado en Proyectos (ABPro) que aplica la derivada a un problema real de ingeniería de software.",
        "Consiste en modelar el costo mensual de alojamiento en la nube en función del número de instancias contratadas.",
        "Mediante la derivada de esa función se determina el punto exacto donde el costo es mínimo.",
        "El resultado se comunica a través de un sitio web con un simulador interactivo.",
    ])

    # ---------- Diapositiva 3: Tema, nociones matemáticas y problemática ----------
    s = slides[2]
    tf = limpiar_y_titular(s.shapes[0].text_frame, "Tema, nociones matemáticas y problemática")
    reposicionar(s.shapes[0], height=Emu(1325563))
    agregar_parrafo(tf, "Tema:", tam=18, negrita=True)
    agregar_bullets(tf, ["Optimización de costos de infraestructura cloud mediante cálculo diferencial."])
    agregar_parrafo(tf, "Nociones matemáticas:", tam=18, negrita=True, espacio_antes=10)
    agregar_bullets(tf, ["Funciones cuadráticas, derivada, puntos críticos y criterio de la segunda derivada."])
    agregar_parrafo(tf, "Problemática:", tam=18, negrita=True, espacio_antes=10)
    agregar_bullets(tf, [
        "Sobredimensionar recursos implica pagar de más; subdimensionarlos afecta el rendimiento del servicio.",
        "Pregunta central: ¿cómo determinar matemáticamente el número de instancias que minimiza el costo?",
    ])

    # ---------- Diapositiva 4: Objetivos y preguntas de investigación ----------
    s = slides[3]
    tf = limpiar_y_titular(s.shapes[0].text_frame, "Objetivos y preguntas de investigación")
    reposicionar(s.shapes[0], height=Emu(1325563))
    agregar_parrafo(tf, "Objetivo general:", tam=18, negrita=True)
    agregar_bullets(tf, ["Modelar el costo de infraestructura cloud como función del número de instancias y aplicar la derivada para hallar el costo mínimo."])
    agregar_parrafo(tf, "Objetivos específicos:", tam=18, negrita=True, espacio_antes=10)
    agregar_bullets(tf, [
        "Construir la función de costo C(x).",
        "Calcular C'(x) e identificar sus puntos críticos.",
        "Verificar el mínimo con el criterio de la segunda derivada.",
        "Desarrollar un simulador interactivo que visualice C(x) y C'(x).",
    ])
    agregar_parrafo(tf, "Pregunta de investigación:", tam=18, negrita=True, espacio_antes=10)
    agregar_bullets(tf, ["¿Cuál es el número óptimo de instancias que minimiza el costo mensual, y cómo se determina usando la derivada?"])

    # ---------- Diapositiva 5: Temas iniciales involucrados ----------
    s = slides[4]
    tf = limpiar_y_titular(s.shapes[0].text_frame, "Temas iniciales involucrados")
    reposicionar(s.shapes[0], height=Emu(1325563))
    agregar_bullets(tf, [
        "Funciones polinomiales (función cuadrática de costo).",
        "Límites y continuidad.",
        "Reglas de derivación.",
        "Puntos críticos y criterio de la primera y segunda derivada.",
        "Optimización aplicada a problemas de ingeniería.",
    ])

    # ---------- Diapositiva 6: Cronograma y materiales ----------
    s = slides[5]
    tf = limpiar_y_titular(s.shapes[0].text_frame, "Cronograma y materiales a utilizar")
    reposicionar(s.shapes[0], height=Emu(1325563))
    # Quitar la imagen de plantilla (tabla vacía) y reemplazarla por una tabla real
    for shape in list(s.shapes):
        if shape.shape_type == 13:  # PICTURE
            shape._element.getparent().remove(shape._element)

    filas_cronograma = [
        ("Investigar problemática y definir tema", "Problemática y tema definidos", "Semana 1-2", "Andrés Carrasco", "Bibliografía de cálculo"),
        ("Formular objetivos y modelo matemático inicial", "Función de costo C(x) planteada", "Semana 3-4", "Andrés Carrasco", "Apuntes de cálculo"),
        ("Aplicar la derivada y validar el mínimo", "Punto crítico verificado con C''(x)", "Semana 5-6", "Andrés Carrasco", "Software de cálculo (GeoGebra/Python)"),
        ("Diseñar y programar el simulador web", "Sitio web + simulador interactivo", "Semana 7-9", "Andrés Carrasco", "HTML/CSS/JS, Plotly.js, Vercel"),
        ("Redactar resultados, conclusiones y bibliografía", "Informe final y presentación", "Semana 10-11", "Andrés Carrasco", "Procesador de texto, fuentes"),
        ("Grabar video y preparar entrega final", "Video + PPT final", "Semana 12", "Andrés Carrasco", "Guion, grabadora de audio, ffmpeg"),
    ]
    n_filas = len(filas_cronograma) + 1
    tabla_shape = s.shapes.add_table(n_filas, 5, IZQ, Emu(1790688), ANCHO, Emu(3600000))
    tabla = tabla_shape.table
    encabezados = ["Tareas", "Productos", "Plazos", "Responsable(s)", "Insumos"]
    for c, texto in enumerate(encabezados):
        celda = tabla.cell(0, c)
        celda.text = texto
        for p in celda.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.name = "Calibri"
    for fila_idx, fila in enumerate(filas_cronograma, start=1):
        for c, valor in enumerate(fila):
            celda = tabla.cell(fila_idx, c)
            celda.text = valor
            for p in celda.text_frame.paragraphs:
                p.font.size = Pt(11)
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = "Calibri"

    # ---------- Diapositiva 7: Justificación de herramientas tecnológicas ----------
    s = slides[6]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Justificación de herramientas tecnológicas")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "Se utilizó un simulador web interactivo (HTML, CSS, JavaScript y la librería Plotly.js) en lugar de solo cálculos manuales.",
        "Permite visualizar en tiempo real la curva de costo C(x) y su derivada C'(x), y modificar parámetros del modelo.",
        "Facilita comprender de forma visual el concepto de punto crítico y optimización, algo difícil de transmitir solo con álgebra.",
        "El sitio se publicó en la nube (Vercel), reforzando la coherencia entre el tema del proyecto (costos cloud) y la herramienta usada para exponerlo.",
    ])

    # ---------- Diapositiva 8: Marco teórico I ----------
    s = slides[7]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Marco teórico I: Funciones y modelo de costo")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "Una función cuadrática C(x) = a·x² + (b+c)·x + d permite representar el costo total en función de la cantidad de recursos (x).",
        "El coeficiente a > 0 refleja que, a partir de cierto punto, agregar más instancias aumenta el costo (sobreaprovisionamiento).",
        "Referencias: Stewart, J. (2018). Cálculo: Trascendentes tempranas; Larson, R. y Edwards, B. (2018). Cálculo.",
    ])

    # ---------- Diapositiva 9: Marco teórico II ----------
    s = slides[8]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Marco teórico II: Derivada y optimización")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "La derivada C'(x) representa la tasa de cambio del costo respecto al número de instancias (costo marginal).",
        "Un punto crítico se obtiene resolviendo C'(x) = 0.",
        "El criterio de la segunda derivada (C''(x) > 0) confirma si ese punto corresponde a un mínimo.",
        "Estos conceptos corresponden a las unidades de límites, derivadas y aplicaciones de la derivada del curso.",
    ])

    # ---------- Diapositiva 10: Marco teórico III ----------
    s = slides[9]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Marco teórico III: Cloud computing y costos")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "Los proveedores cloud (AWS, Azure, Vercel) cobran según los recursos utilizados (instancias, cómputo, almacenamiento).",
        "Un mal dimensionamiento de recursos es un problema real de la industria: sobrecostos o baja disponibilidad del servicio.",
        "Fuentes: documentación de precios de AWS y Azure, y buenas prácticas de \"cost optimization\" en la nube.",
    ])

    # ---------- Diapositiva 11: Desarrollo, revisión de ideas y producto ----------
    s = slides[10]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Desarrollo, revisión de ideas y producto")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "Se evaluaron distintos modelos de costo (lineal y cuadrático), optando por un modelo cuadrático simplificado que resume infraestructura + operación.",
        "Se derivó el modelo, se resolvió C'(x) = 0 y se verificó el mínimo con la segunda derivada.",
        "Se decidió comunicar el resultado mediante un producto tecnológico -un sitio web con simulador interactivo- en vez de solo un informe estático.",
        "El simulador permite ingresar distintos parámetros (a, b+c, d) y observar cómo cambia el punto óptimo.",
    ])

    # ---------- Diapositiva 12: Presentación de producto o prototipo ----------
    s = slides[11]
    caja_titulo = s.shapes[0]  # "Diapositiva 12: Presentación de producto o prototipo"
    limpiar_y_titular(caja_titulo.text_frame, "Presentación de producto o prototipo", tam_titulo=24)
    # NO tocar s.shapes[1] (contiene el link ya agregado por agregar_link_ppt.py)
    link_shape = s.shapes[1]
    top_contenido = Emu(link_shape.top + link_shape.height + 150000)
    caja_contenido = s.shapes.add_textbox(link_shape.left, top_contenido, Emu(9500000), Emu(2500000))
    tf_contenido = caja_contenido.text_frame
    tf_contenido.word_wrap = True
    agregar_bullets(tf_contenido, [
        "Producto: sitio web (HTML, CSS, JavaScript) con un simulador interactivo construido con Plotly.js.",
        "Grafica en tiempo real la función de costo C(x) y su derivada C'(x), marcando el punto donde C'(x) = 0.",
        "Publicado en la nube mediante Vercel, con despliegue continuo desde GitHub.",
    ])
    # el primer párrafo quedó vacío por defecto al crear el textbox; limpiarlo
    if tf_contenido.paragraphs[0].runs == []:
        tf_contenido.paragraphs[0]._p.getparent().remove(tf_contenido.paragraphs[0]._p)

    # ---------- Diapositiva 13: Resultados y conclusiones ----------
    s = slides[12]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Resultados y conclusiones")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_parrafo(tf, "Resultados:", tam=18, negrita=True)
    agregar_bullets(tf, [
        "Modelo: C(x) = 0.5x² − 20x + 500  ;  C'(x) = x − 20",
        "Punto crítico: x = 20  ;  C''(x) = 1 > 0 → mínimo",
        "Costo mínimo: C(20) = 300 USD/mes",
        "Comparación: 5 instancias → 412,5 USD ; 20 instancias (óptimo) → 300 USD ; 35 instancias → 412,5 USD (ahorro ≈ 27%)",
    ])
    agregar_parrafo(tf, "Conclusión:", tam=18, negrita=True, espacio_antes=10)
    agregar_bullets(tf, [
        "La derivada permite tomar decisiones de arquitectura de software basadas en evidencia matemática, y no en estimaciones arbitrarias.",
    ])

    # ---------- Diapositiva 14: Bibliografía y fuentes de información ----------
    s = slides[13]
    caja = s.shapes[0]
    tf = limpiar_y_titular(caja.text_frame, "Bibliografía y fuentes de información")
    reposicionar(caja, height=TOPE_INFERIOR - Emu(365125))
    agregar_bullets(tf, [
        "Stewart, J. (2018). Cálculo de una variable: Trascendentes tempranas (8ª ed.). Cengage Learning.",
        "Larson, R., & Edwards, B. (2018). Cálculo (10ª ed.). McGraw-Hill.",
        "Amazon Web Services. (2024). Cloud economics and pricing. https://aws.amazon.com/economics/",
        "Microsoft Azure. (2024). Pricing calculator y documentación de cost management. https://azure.microsoft.com/pricing/",
        "Vercel Inc. (2024). Vercel pricing documentation. https://vercel.com/pricing",
        f"Sitio del proyecto: {URL_SITIO}",
    ], tam=16)

    prs.save(ARCHIVO)
    print(f"Contenido completado en {ARCHIVO}")


if __name__ == "__main__":
    main()
