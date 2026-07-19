"""Aplica al PPT una identidad visual coherente con el sitio web del proyecto
(mismos colores y jerarquía tipográfica que web/styles.css), sin depender de
fuentes que no vienen instaladas por defecto en Windows/Mac (se usa Calibri,
que ya trae la plantilla, evitando saltos de fuente al abrir el archivo en
otro equipo).

Se ejecuta DESPUÉS de completar_contenido_ppt.py: reescribe fondo, colores de
texto, viñetas reales (no texto "• ") y la tabla del cronograma, además de
actualizar el cronograma con las fechas reales de Contenido_Proyecto.md.

Uso:
    python3 estilizar_ppt.py
"""
import os
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ARCHIVO = os.path.join(os.path.dirname(__file__), "..", "Presentación Proyecto.pptx")

# Paleta tomada de web/styles.css (:root)
BG = RGBColor(0x05, 0x06, 0x11)
TEXTO = RGBColor(0xE6, 0xEB, 0xF5)
TEXTO_DIM = RGBColor(0x94, 0xA3, 0xB8)
CIAN = RGBColor(0x22, 0xD3, 0xEE)
VIOLETA = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x5E, 0xEA, 0xD4)
BORDE = RGBColor(0x2B, 0x30, 0x45)
SURFACE = RGBColor(0x0D, 0x0F, 0x1E)

FUENTE = "Calibri"


def fondo_oscuro(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def barra_superior(slide, prs):
    """Línea de acento arriba de cada diapositiva, degradado cian→violeta,
    igual al gradiente de marca del sitio (--gradient-brand)."""
    from pptx.enum.shapes import MSO_SHAPE
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Emu(45720))
    barra.line.fill.background()
    barra.shadow.inherit = False
    fill = barra.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = CIAN
    stops[0].position = 0.0
    stops[1].color.rgb = VIOLETA
    stops[1].position = 1.0
    fill.gradient_angle = 0.0
    return barra


def set_bullet_char(paragraph, color=TEAL, char="\u25aa", size_pct=80):
    """Reemplaza el bullet por defecto por uno con color propio (nivel 0)."""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buSzPct", "a:buFont"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    buClr.append(srgb)
    buSzPct = pPr.makeelement(qn("a:buSzPct"), {"val": str(size_pct * 1000)})
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buClr)
    pPr.append(buSzPct)
    pPr.append(buFont)
    pPr.append(buChar)
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")


def quitar_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    buNone = pPr.makeelement(qn("a:buNone"), {})
    pPr.append(buNone)


def estilizar_parrafo(paragraph, es_encabezado_seccion):
    texto = paragraph.text
    if texto.startswith("•  "):
        # quitar el "• " manual y usar un bullet real de PowerPoint,
        # preservando el tamaño de fuente que ya traía el párrafo
        nuevo_texto = texto[3:]
        tam_original = None
        if paragraph.runs and paragraph.runs[0].font.size:
            tam_original = paragraph.runs[0].font.size
        for run in list(paragraph.runs):
            run._r.getparent().remove(run._r)
        run = paragraph.add_run()
        run.text = nuevo_texto
        run.font.size = tam_original or Pt(16)
        run.font.color.rgb = TEXTO
        run.font.name = FUENTE
        set_bullet_char(paragraph)
        paragraph.line_spacing = 1.15
        paragraph.space_after = Pt(6)
    else:
        quitar_bullet(paragraph)
        for run in paragraph.runs:
            run.font.name = FUENTE
            if run.font.bold:
                run.font.color.rgb = VIOLETA
            else:
                run.font.color.rgb = TEXTO
        paragraph.space_after = Pt(4)
        if es_encabezado_seccion:
            paragraph.space_before = Pt(10)


def estilizar_caja_texto(shape, es_titulo_slide=False):
    tf = shape.text_frame
    for i, p in enumerate(tf.paragraphs):
        if not p.runs:
            continue
        if i == 0 and es_titulo_slide:
            for run in p.runs:
                run.font.color.rgb = CIAN
                run.font.bold = True
                run.font.name = FUENTE
            quitar_bullet(p)
            p.space_after = Pt(14)
        else:
            estilizar_parrafo(p, es_encabezado_seccion=(p.runs[0].font.bold if p.runs else False))


def estilizar_tabla_cronograma(slide):
    tabla_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            tabla_shape = shape
            break
    if tabla_shape is None:
        return
    tabla = tabla_shape.table

    # Fechas reales (coinciden con Contenido_Proyecto.md)
    filas_reales = [
        ("Definir problema y objetivos", "Problemática y objetivos definidos", "7 julio", "Andrés Carrasco", "Apuntes de Cálculo Unidad 3"),
        ("Investigar costos reales de proveedores cloud y modelar C(x)", "Función de costo C(x) planteada", "9 julio", "Andrés Carrasco", "Documentación de precios AWS/Azure"),
        ("Aplicar la derivada, hallar el punto óptimo y construir el simulador", "C'(x), punto crítico x=20 y simulador web", "14 julio", "Andrés Carrasco", "Plotly.js, HTML/CSS/JS"),
        ("Elaborar PPT y grabar el video de presentación", "PPT + video con guion", "14-20 julio", "Andrés Carrasco", "Guion, grabadora de audio, ffmpeg"),
        ("Revisión final y entrega", "Proyecto completo entregado", "21 julio", "Andrés Carrasco", "Sitio web publicado en Vercel"),
    ]
    # quitar filas de más si el nuevo contenido tiene menos filas que la tabla original
    filas_actuales = len(tabla.rows) - 1
    filas_necesarias = len(filas_reales)
    tbl = tabla_shape.table._tbl
    while len(tabla.rows) - 1 > filas_necesarias:
        tr = tbl.tr_lst[-1]
        tbl.remove(tr)

    encabezados = ["Tareas", "Productos", "Plazos", "Responsable(s)", "Insumos"]
    for c, texto in enumerate(encabezados):
        celda = tabla.cell(0, c)
        celda.text = texto
        celda.fill.solid()
        celda.fill.fore_color.rgb = SURFACE
        for p in celda.text_frame.paragraphs:
            quitar_bullet(p)
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.name = FUENTE
                r.font.color.rgb = CIAN
    for fila_idx, fila in enumerate(filas_reales, start=1):
        for c, valor in enumerate(fila):
            celda = tabla.cell(fila_idx, c)
            celda.text = valor
            celda.fill.solid()
            celda.fill.fore_color.rgb = BG
            for p in celda.text_frame.paragraphs:
                quitar_bullet(p)
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = FUENTE
                    r.font.color.rgb = TEXTO_DIM if c != 2 else TEAL


def main():
    prs = Presentation(ARCHIVO)

    for idx, slide in enumerate(prs.slides):
        fondo_oscuro(slide)
        barra_superior(slide, prs)

        # snapshot de ids ANTES de tocar nada: slide.shapes[0] crea un objeto
        # proxy nuevo en cada acceso, por lo que comparar con "is" nunca es
        # verdadero; comparamos por shape_id en su lugar.
        primer_shape_id = list(slide.shapes)[0].shape_id
        for shape in list(slide.shapes):
            if shape.has_table:
                continue
            if not shape.has_text_frame:
                continue
            es_titulo = shape.shape_id == primer_shape_id
            estilizar_caja_texto(shape, es_titulo_slide=es_titulo)

        if idx == 5:  # Diapositiva 6: Cronograma
            estilizar_tabla_cronograma(slide)

        if idx == 0:  # Portada: subir un poco el contraste del subtítulo
            for shape in slide.shapes:
                if shape.has_text_frame and "Andrés" in shape.text_frame.text:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = TEXTO_DIM

    prs.save(ARCHIVO)
    print(f"Estilo aplicado en {ARCHIVO}")


if __name__ == "__main__":
    main()
